from pptx import Presentation

def extract_ppt(file):
    prs = Presentation(file)
    notes_text = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text = notes_slide.notes_text_frame.text
            notes_text.append(text)

    return notes_text
